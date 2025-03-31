import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches
from sklearn.model_selection import train_test_split
from sklearn.ensemble import RandomForestRegressor
from sklearn.metrics import mean_absolute_error, mean_squared_error
from statsmodels.tsa.seasonal import seasonal_decompose
import warnings
warnings.filterwarnings("ignore")

# Load dataset
df = pd.read_csv(r"C:\Users\theba\OneDrive\Desktop\Projects\Project 9 Coffee Sale Data\index.csv", parse_dates=['date'])

# Data Cleaning
df.columns = df.columns.str.strip()  # Strip spaces from column names
df.rename(columns={'money': 'sales'}, inplace=True)  # Rename 'money' to 'sales'
df = df.infer_objects()  # Convert object columns to numeric or datetime if possible
df = df.interpolate(method='linear')  # Interpolate missing values

# Feature Engineering
df['sales_MA_7'] = df['sales'].rolling(window=7).mean()  # 7-day moving average
df['sales_MA_30'] = df['sales'].rolling(window=30).mean()  # 30-day moving average
df['day_of_week'] = df['date'].dt.dayofweek
df['month'] = df['date'].dt.month

# Seasonal Decomposition
decomposition = seasonal_decompose(df['sales'], model='additive', period=30)
fig = decomposition.plot()
fig.set_size_inches(10, 6)
plt.savefig('seasonal_decomposition.png')
plt.close()

# Line Plot for Daily Sales
plt.figure(figsize=(10, 6))
sns.lineplot(data=df, x='date', y='sales')
plt.title('Daily Coffee Sales')
plt.xlabel('Date')
plt.ylabel('Sales')
plt.savefig('daily_sales.png')
plt.close()

# Moving Averages Plot
plt.figure(figsize=(10, 6))
plt.plot(df['date'], df['sales'], label='Daily Sales')
plt.plot(df['date'], df['sales_MA_7'], label='7-Day Moving Average', linestyle='--')
plt.plot(df['date'], df['sales_MA_30'], label='30-Day Moving Average', linestyle='--')
plt.title('Daily Sales and Moving Averages')
plt.xlabel('Date')
plt.ylabel('Sales')
plt.legend()
plt.savefig('moving_averages.png')
plt.close()

# Heatmap for Feature Correlation
numeric_columns = df.select_dtypes(include=['number']).columns  # Select numeric columns
correlation_data = df[numeric_columns].corr()  # Correlation matrix
plt.figure(figsize=(8, 6))
sns.heatmap(correlation_data, annot=True, cmap='coolwarm', fmt='.2f')
plt.title('Feature Correlation Heatmap')
plt.savefig('correlation_heatmap.png')
plt.close()

# Machine Learning: Random Forest
df['sales_lag_1'] = df['sales'].shift(1)  # Lag feature for sales
df.dropna(inplace=True)  # Drop rows with NaN values from lag and moving averages
X = df[['day_of_week', 'month', 'sales_lag_1', 'sales_MA_7']]
y = df['sales']

X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)

model_rf = RandomForestRegressor(random_state=42)
model_rf.fit(X_train, y_train)
y_pred_rf = model_rf.predict(X_test)

mae_rf = mean_absolute_error(y_test, y_pred_rf)
mse_rf = mean_squared_error(y_test, y_pred_rf)
rmse_rf = mse_rf ** 0.5

# Create PowerPoint Presentation
presentation = Presentation()

# Slide 1: Title Slide
slide_1 = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Coffee Sales Data Analysis and Forecasting"
subtitle.text = "Data Cleaning, Machine Learning, and Insights"

# Slide 2: Dataset Overview
slide_2 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_2.shapes.title
content = slide_2.placeholders[1]
title.text = "Dataset Overview"
content.text = ("The dataset contains daily sales data for a coffee chain.\n"
                "Key features include sales, date, and seasonality patterns.")

# Slide 3: Data Cleaning
slide_3 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_3.shapes.title
content = slide_3.placeholders[1]
title.text = "Data Cleaning"
content.text = ("- Stripped spaces from column names.\n"
                "- Renamed 'money' column to 'sales'.\n"
                "- Interpolated missing values using linear interpolation.")

# Slide 4: Daily Sales Plot
slide_4 = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide_4.shapes.title
title.text = "Daily Coffee Sales"
slide_4.shapes.add_picture('daily_sales.png', Inches(1), Inches(1), width=Inches(8), height=Inches(5))

# Slide 5: Moving Averages Plot
slide_5 = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide_5.shapes.title
title.text = "Daily Sales and Moving Averages"
slide_5.shapes.add_picture('moving_averages.png', Inches(1), Inches(1), width=Inches(8), height=Inches(5))

# Slide 6: Seasonal Decomposition
slide_6 = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide_6.shapes.title
title.text = "Seasonal Decomposition"
slide_6.shapes.add_picture('seasonal_decomposition.png', Inches(1), Inches(1), width=Inches(8), height=Inches(5))

# Slide 7: Correlation Heatmap
slide_7 = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide_7.shapes.title
title.text = "Feature Correlation Heatmap"
slide_7.shapes.add_picture('correlation_heatmap.png', Inches(1), Inches(1), width=Inches(8), height=Inches(5))

# Slide 8: Random Forest Model Evaluation
slide_8 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_8.shapes.title
content = slide_8.placeholders[1]
title.text = "Model Evaluation"
content.text = (f"Random Forest Regressor:\n"
                f"- Mean Absolute Error: {mae_rf:.2f}\n"
                f"- Mean Squared Error: {mse_rf:.2f}\n"
                f"- Root Mean Squared Error: {rmse_rf:.2f}")

# Save the PowerPoint Presentation
file_name = "Coffee_Sales_Analysis_Presentation.pptx"
presentation.save(file_name)

print(f"Presentation saved as {file_name}")
